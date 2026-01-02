#!/usr/bin/env python3
"""Script to update the Autho.R presentation with product information and logo."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import re

# Load the presentation - use original template
pptx_path = "/Users/_malcolmadams/projects/author-rpa/Blue Modern Robotic Process Automation for Business Presentation.pptx"
prs = Presentation(pptx_path)

# Logo path - use white version for dark backgrounds
LOGO_PATH = "/Users/_malcolmadams/projects/author-rpa/process-automation-white.png"

# Logo size settings (aesthetically pleasing - matches screenshot)
LOGO_HEIGHT = Inches(1.2)  # Larger, prominent size matching the screenshot

# Product information
PRODUCT_NAME = "Autho.R"
TAGLINE = "Enterprise-Grade Robotic Process Automation"
COMPANY = "OnticWorks.io"
WEBSITE = "onticworks.io"
SALES_EMAIL = "sales@onticworks.io"
SUPPORT_EMAIL = "support@onticworks.io"
APP_URL = "author-rpa-app.azurewebsites.net"

# Improved content for different slides with better flow
SLIDE_CONTENT = {
    2: {  # Introduction to Autho.R
        "main": "Autho.R is a powerful Python-based automation platform that transforms how businesses handle repetitive tasks. From document generation to web scraping, database operations to API integrations—Autho.R handles it all. Available on Azure and IBM Cloud Marketplaces.",
    },
    3: {  # Why Autho.R Matters
        "main": "In today's fast-paced business environment, manual processes drain resources and introduce errors. Autho.R eliminates these bottlenecks by automating routine tasks, allowing your team to focus on strategic initiatives that drive growth.",
    },
    4: {  # Key Features
        "main": "Nine powerful modules working together: Create documents (Word, PDF, Markdown), process spreadsheets (Excel, CSV), scrape websites, integrate APIs, manage databases (SQLite, PostgreSQL, MySQL), automate desktop tasks, send emails, and orchestrate complex workflows.",
    },
    5: {  # Flexible Pricing
        "main": "Choose the plan that fits your needs. Start free with 100 tasks/month, or scale to Enterprise for unlimited automation. All plans include our intuitive web UI, Python API access, and Azure/IBM Cloud deployment options.",
    },
    6: {  # Real-World Applications
        "main": "Autho.R powers diverse automation scenarios: automated report generation, invoice processing, data migration, competitive monitoring via web scraping, CRM synchronization, document templating, and multi-system integrations—all without writing complex code.",
    },
    7: {  # Getting Started
        "step1": "Select Your Plan: Choose from Free, Starter ($29/mo), Professional ($79/mo), Business ($199/mo), or Enterprise ($499+/mo) based on your task volume.",
        "step2": "Deploy in Minutes: One-click deployment from Azure Marketplace or IBM Cloud. No infrastructure setup required.",
        "step3": "Build Your Workflows: Use our intuitive web interface or Python API to create automation workflows that connect your systems.",
        "step4": "Scale With Confidence: Monitor performance, optimize processes, and scale seamlessly as your automation needs grow.",
    },
    8: {  # Success Stories
        "example1": "Document Automation: A financial services firm reduced report generation time from 4 hours to 15 minutes, processing 500+ client reports daily with zero errors.",
        "example2": "Data Integration: An e-commerce company automated inventory sync across 3 platforms, eliminating manual entry and reducing stock discrepancies by 95%.",
    },
    9: {  # The Future of Work
        "main": "Automation isn't just about efficiency—it's about empowering your team. Autho.R handles the repetitive work so your people can focus on creativity, strategy, and customer relationships.",
        "future": "Join forward-thinking organizations already transforming their operations with Autho.R. Start your automation journey today with our free tier.",
    },
    10: {  # Get Started Today
        "main": f"Ready to automate? Visit {APP_URL} to start free, or contact {SALES_EMAIL} for enterprise solutions. Our team is here to help you succeed.",
    },
}

def is_lorem_ipsum(text):
    """Check if text contains Lorem ipsum or similar placeholder text."""
    lorem_indicators = [
        'lorem ipsum', 'dolor sit amet', 'consectetur adipiscing',
        'elit', 'proin', 'nulla', 'adipiscing', 'tempor incididunt',
        'labore et dolore', 'magna aliqua', 'enim ad minim', 'veniam',
        'quis nostrud', 'exercitation', 'ullamco', 'laboris',
        'consequat', 'duis aute', 'irure dolor', 'reprehenderit',
        'voluptate', 'velit esse', 'cillum', 'fugiat nulla', 'pariatur',
        'excepteur sint', 'occaecat', 'cupidatat', 'proident', 'sunt in culpa',
        'officia deserunt', 'mollit anim', 'est laborum', 'euismod',
        'aliquip ex ea', 'commodo consequat'
    ]
    text_lower = text.lower()
    return any(indicator in text_lower for indicator in lorem_indicators)

def get_replacement_text(slide_num, text_index):
    """Get appropriate replacement text based on slide and position."""
    content = SLIDE_CONTENT.get(slide_num, {})

    if slide_num == 7:  # Implementation steps
        steps = ["step1", "step2", "step3", "step4"]
        if text_index < len(steps):
            return content.get(steps[text_index], content.get("main", ""))
    elif slide_num == 8:  # Real-world examples
        examples = ["example1", "example2"]
        if text_index < len(examples):
            return content.get(examples[text_index], content.get("main", ""))
    elif slide_num == 9:  # Conclusion
        texts = ["main", "future"]
        if text_index < len(texts):
            return content.get(texts[text_index], content.get("main", ""))

    return content.get("main", "Autho.R - Enterprise-Grade Robotic Process Automation by OnticWorks.io")

def get_group_text(shape):
    """Get all text from a group shape."""
    text = ""
    if hasattr(shape, 'shapes'):
        for s in shape.shapes:
            if hasattr(s, 'text_frame'):
                text += s.text_frame.text + " "
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                text += get_group_text(s)
    return text.strip()

def replace_logo_groups_and_pictures(slide, logo_path):
    """Find and replace logo groups (containing OnticWorks.io or company text) and pictures, or add logo if none found."""
    shapes_to_remove = []
    logo_positions = []

    for shape in slide.shapes:
        try:
            shape_type = shape.shape_type
            should_replace = False

            # Check if it's a picture - always replace
            if shape_type == MSO_SHAPE_TYPE.PICTURE or shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE:
                should_replace = True
                print(f"    Found PICTURE: {shape.name}")

            # Check if it's a group that might contain the company logo
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                group_text = get_group_text(shape).lower()
                # Check if this group contains company name text (the logo group)
                if 'onticworks' in group_text or 'thynk' in group_text:
                    should_replace = True
                    print(f"    Found logo GROUP: {shape.name} with text: {group_text[:50]}")
                # Also check if it's in the top-left logo position (small group in corner)
                elif shape.top < Inches(1) and shape.left < Inches(2) and shape.width < Inches(2):
                    should_replace = True
                    print(f"    Found corner GROUP (likely logo): {shape.name}")

            if should_replace:
                logo_positions.append({
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'name': shape.name
                })
                shapes_to_remove.append(shape)

        except Exception as e:
            print(f"    Error checking shape {shape.name}: {e}")

    print(f"    Found {len(shapes_to_remove)} shapes to replace")

    # Remove old shapes
    for shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
            print(f"    Removed: {shape.name}")
        except Exception as e:
            print(f"    Could not remove {shape.name}: {e}")

    # Add new logo - either at found positions or at default top-left position
    if logo_positions:
        for pos in logo_positions:
            try:
                slide.shapes.add_picture(
                    logo_path,
                    pos['left'],
                    pos['top'],
                    height=LOGO_HEIGHT
                )
                print(f"    Added new logo replacing: {pos['name']}")
            except Exception as e:
                print(f"    Could not add logo: {e}")
    else:
        # No existing logo found - add at top-left corner
        try:
            slide.shapes.add_picture(
                logo_path,
                Inches(0.3),   # Left margin
                Inches(0.25),  # Top margin
                height=LOGO_HEIGHT
            )
            print(f"    Added new logo at top-left corner")
        except Exception as e:
            print(f"    Could not add logo: {e}")

    return len(shapes_to_remove) if shapes_to_remove else 1  # Return 1 if we added a new logo

def update_presentation():
    """Update the presentation with Autho.R content and logo."""

    print(f"Loaded presentation with {len(prs.slides)} slides\n")

    # Common text replacements
    common_replacements = [
        ("Thynk Unlimited", COMPANY),
        ("OnticWorks.io", COMPANY),
        ("123-456-7890", WEBSITE),
        ("hello@reallygreatsite.com", SALES_EMAIL),
        ("123 Anywhere St., Any City, ST 12345", APP_URL),
        ("Robotic Process Automation", PRODUCT_NAME),
        ("Introduction to Autho.R", f"Introduction to {PRODUCT_NAME}"),
        ("FOR BUSINESS", f"by {COMPANY}"),
        ("Transforming Business Efficiency with Intelligent Automation", TAGLINE),
        ("Why RPA is", "Why Autho.R is"),
        ("Why Autho.R is Important for Businesses", "Why Automation Matters"),
        ("Key Features of RPA", f"Key Features of {PRODUCT_NAME}"),
        ("Key Features of Autho.R", f"What {PRODUCT_NAME} Can Do"),
        ("Benefits of Implementing RPA", "Flexible Pricing for Every Need"),
        ("Benefits of Implementing Autho.R", "Flexible Pricing for Every Need"),
        ("Applications of RPA in Business", "Real-World Applications"),
        ("Applications of Autho.R in Business", "Real-World Applications"),
        ("Steps to Implement RPA in Your Business", "Getting Started is Easy"),
        ("Steps to Implement Autho.R in Your Business", "Getting Started is Easy"),
        ("Real-World Examples of RPA Success", "Success Stories"),
        ("Real-World Examples of Autho.R Success", "Success Stories"),
        ("Conclusion and Future Outlook", "The Future of Work is Automated"),
    ]

    # Process each slide
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n{'='*50}")
        print(f"Processing Slide {slide_num}...")
        print(f"{'='*50}")
        lorem_count = 0

        # Replace logos (both pictures and logo groups)
        try:
            replaced = replace_logo_groups_and_pictures(slide, LOGO_PATH)
            print(f"  Replaced {replaced} logo element(s) on slide {slide_num}")
        except Exception as e:
            print(f"  Error replacing logos on slide {slide_num}: {e}")

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for old, new in common_replacements:
                            if old.lower() in run.text.lower():
                                pattern = re.compile(re.escape(old), re.IGNORECASE)
                                run.text = pattern.sub(new, run.text)

                        if is_lorem_ipsum(run.text) and len(run.text) > 20:
                            replacement = get_replacement_text(slide_num, lorem_count)
                            print(f"  Replacing placeholder text with: {replacement[:50]}...")
                            run.text = replacement
                            lorem_count += 1

            if hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for old, new in common_replacements:
                                    if old.lower() in run.text.lower():
                                        pattern = re.compile(re.escape(old), re.IGNORECASE)
                                        run.text = pattern.sub(new, run.text)

                                if is_lorem_ipsum(run.text) and len(run.text) > 20:
                                    replacement = get_replacement_text(slide_num, lorem_count)
                                    run.text = replacement
                                    lorem_count += 1

    # Save the updated presentation to new file
    output_path = "/Users/_malcolmadams/projects/author-rpa/Autho.R_Presentation.pptx"
    prs.save(output_path)
    print(f"\n\nSaved updated presentation to: {output_path}")
    print("Original template preserved.")

    return output_path

if __name__ == "__main__":
    update_presentation()
